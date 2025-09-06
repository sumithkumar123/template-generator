import os
from openai import OpenAI
import docx
import fitz  # PyMuPDF
import pandas as pd
from openpyxl import load_workbook
from pptx import Presentation
from PIL import Image
import json

# --- Configuration ---
# IMPORTANT: Set your OpenAI API key in your environment variables.
# For a hackathon, you can temporarily hardcode it here, but it's not recommended.
# client = OpenAI(api_key="YOUR_OPENAI_API_KEY")
client = OpenAI(api_key=os.environ.get("OPENAI_API_KEY"))
UPLOAD_DIRECTORY = "uploads"


def parse_document(filename: str) -> str:
    """
    Parses the content of an uploaded file (PDF, DOCX, XLSX, PPTX).
    Adds citation markers for page/paragraph/slide/sheet numbers.
    """
    filepath = os.path.join(UPLOAD_DIRECTORY, filename)
    content_parts = []
    
    try:
        if filename.lower().endswith(".pdf"):
            doc = fitz.open(filepath)
            for page_num, page in enumerate(doc):
                text = page.get_text("text")
                if text.strip():
                    content_parts.append(f"[START PAGE {page_num + 1}]\n{text}\n[END PAGE {page_num + 1}]")
            doc.close()
        
        elif filename.lower().endswith(".docx"):
            doc = docx.Document(filepath)
            for para_num, para in enumerate(doc.paragraphs):
                if para.text.strip():
                    content_parts.append(f"[START PARA {para_num + 1}]\n{para.text}\n[END PARA {para_num + 1}]")
        
        elif filename.lower().endswith((".xlsx", ".xls")):
            # Parse Excel files
            workbook = load_workbook(filepath, data_only=True)
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                content_parts.append(f"[START SHEET: {sheet_name}]")
                
                # Extract data from sheet
                sheet_data = []
                for row_num, row in enumerate(sheet.iter_rows(values_only=True), 1):
                    if any(cell is not None for cell in row):
                        row_data = [str(cell) if cell is not None else "" for cell in row]
                        sheet_data.append(f"Row {row_num}: {' | '.join(row_data)}")
                
                if sheet_data:
                    content_parts.append("\n".join(sheet_data))
                content_parts.append(f"[END SHEET: {sheet_name}]")
        
        elif filename.lower().endswith(".pptx"):
            # Parse PowerPoint files
            presentation = Presentation(filepath)
            for slide_num, slide in enumerate(presentation.slides, 1):
                content_parts.append(f"[START SLIDE {slide_num}]")
                
                # Extract text from slide
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                if slide_text:
                    content_parts.append("\n".join(slide_text))
                content_parts.append(f"[END SLIDE {slide_num}]")

        else:
            # Simple text file fallback
            with open(filepath, 'r', encoding='utf-8') as f:
                content_parts.append(f.read())
                
    except Exception as e:
        print(f"Error parsing {filename}: {e}")
        return f"Error parsing file: {filename}"

    return "\n".join(content_parts)


def generate_demo_content(section_title: str, context: str) -> str:
    """
    Generates demo content when OpenAI API is unavailable.
    This allows users to see the application functionality working.
    """
    # Extract filenames from context for citations
    filenames = []
    for line in context.split('\n'):
        if line.startswith('[START DOCUMENT:'):
            filename = line.replace('[START DOCUMENT: ', '').replace(']', '')
            filenames.append(filename)
    
    # Generate appropriate demo content based on section title
    section_lower = section_title.lower()
    
    if 'executive' in section_lower or 'summary' in section_lower:
        return f"""Based on analysis of the uploaded documents, this {section_title.lower()} provides key insights and strategic recommendations.

**Key Highlights:**
- Comprehensive analysis reveals significant opportunities for improvement
- Data-driven insights support strategic decision making
- Implementation recommendations align with organizational objectives

The findings indicate strong performance indicators across multiple metrics [cite: {filenames[0] if filenames else 'document'}, page 1]. Strategic initiatives should focus on leveraging these insights for maximum organizational impact.

*Note: This is demo content. Connect a valid OpenAI API key for AI-generated analysis.*"""
    
    elif 'finding' in section_lower or 'result' in section_lower:
        return f"""The analysis reveals several critical findings that warrant immediate attention and strategic action.

**Primary Findings:**

1. **Performance Metrics:** Analysis demonstrates strong baseline performance with opportunities for optimization [cite: {filenames[0] if filenames else 'document'}, page 2]

2. **Operational Efficiency:** Current processes show potential for 15-25% efficiency improvements through targeted interventions

3. **Risk Assessment:** Identified key risk factors that require mitigation strategies and ongoing monitoring

**Supporting Evidence:**
The data consistently supports these conclusions across multiple evaluation criteria [cite: {filenames[0] if filenames else 'document'}, page 3].

*Note: This is demo content. Connect a valid OpenAI API key for AI-generated analysis.*"""
    
    elif 'conclusion' in section_lower or 'recommendation' in section_lower:
        return f"""Based on comprehensive analysis of the source materials, the following conclusions and recommendations emerge:

**Strategic Recommendations:**

1. **Immediate Actions:** Implement high-priority initiatives within the next 90 days
2. **Medium-term Goals:** Establish sustainable processes for ongoing improvement
3. **Long-term Vision:** Align strategic objectives with organizational mission

**Implementation Roadmap:**
- Phase 1: Foundation building and stakeholder alignment
- Phase 2: Pilot program execution and feedback integration  
- Phase 3: Full-scale implementation and performance monitoring

The evidence strongly supports these recommendations as the optimal path forward [cite: {filenames[0] if filenames else 'document'}, page 4].

*Note: This is demo content. Connect a valid OpenAI API key for AI-generated analysis.*"""
    
    else:
        return f"""This section presents a detailed analysis of {section_title.lower()} based on the provided source materials.

**Overview:**
The comprehensive review reveals important insights that inform strategic decision-making and operational planning.

**Key Points:**
- Data analysis supports evidence-based conclusions
- Multiple sources validate the primary findings
- Recommendations align with best practices and industry standards

**Analysis:**
The evidence demonstrates clear patterns and trends that support the following observations [cite: {filenames[0] if filenames else 'document'}, page 1]. These findings provide a solid foundation for informed decision-making and strategic planning.

*Note: This is demo content. Connect a valid OpenAI API key for AI-generated analysis.*"""


def enhanced_generate_report_section(section_title: str, context: str, tone: str = "professional", style: str = "analytical") -> str:
    """
    Enhanced version that supports tone and style customization.
    Calls the AI model to generate content for a specific section of the template.
    """
    system_prompt = f"""
    You are a world-class business consultant and your task is to generate a section of a report.
    - You will be given a section title and the full context from source documents.
    - Your response MUST be based ONLY on the provided context.
    - You MUST include evidence-backed citations after every statement or claim.
    - Citations should reference the source, like this: [cite: filename, page X], [cite: filename, slide X], [cite: filename, sheet: SheetName], or [cite: filename, para X].
    - The response should be in well-structured Markdown format.
    - Use a {tone} tone and {style} writing style.
    - Be comprehensive and thorough in your analysis.
    """
    
    user_prompt = f"""
    Generate the content for the section: "{section_title}".

    Here is the full context from the source documents:
    ---
    {context}
    ---
    """
    
    try:
        response = client.chat.completions.create(
            model="gpt-4o",  # Use a powerful model for high-quality output
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt},
            ]
        )
        return response.choices[0].message.content
    except Exception as e:
        print(f"Error calling OpenAI API: {e}")
        # Fallback to demo content when API fails
        return generate_demo_content(section_title, context)


def generate_report_section(section_title: str, context: str) -> str:
    """
    Backwards compatible version that uses default tone and style.
    """
    return enhanced_generate_report_section(section_title, context, "professional", "analytical")


def create_final_document(template: list[str], filenames: list[str], tone: str = "professional", style: str = "analytical") -> str:
    """
    Enhanced version with tone and style support.
    Orchestrates the document creation process.
    1. Parses all source files to create a combined context.
    2. Iterates through the template sections, generating content for each.
    3. Assembles the final report.
    """
    # 1. Create Combined Context
    full_context = []
    for filename in filenames:
        parsed_text = parse_document(filename)
        full_context.append(f"[START DOCUMENT: {filename}]\n{parsed_text}\n[END DOCUMENT: {filename}]")
    
    combined_context = "\n\n".join(full_context)
    
    # 2. Generate Each Section
    document_parts = []
    for section in template:
        # Add section title
        document_parts.append(f"## {section}\n")
        
        # Generate section content with tone and style
        section_content = enhanced_generate_report_section(section, combined_context, tone, style)
        document_parts.append(section_content)
        document_parts.append("\n---\n") # Separator
        
    return "\n".join(document_parts)


def ask_clarifying_questions(template: list[str], filenames: list[str]) -> str:
    """
    AI-powered function to ask clarifying questions based on template and uploaded files.
    """
    files_summary = []
    for filename in filenames:
        parsed_content = parse_document(filename)[:500]  # Get first 500 chars for summary
        files_summary.append(f"- {filename}: {parsed_content[:100]}...")
    
    files_text = "\n".join(files_summary)
    template_text = ", ".join(template)
    
    system_prompt = """
    You are an AI consultant assistant. Based on the user's template sections and uploaded files, 
    generate 2-3 clarifying questions that would help create a better, more targeted document.
    Focus on:
    - Missing information that would be valuable
    - Scope clarification
    - Audience and purpose questions
    - Specific requirements or constraints
    
    Keep questions concise and actionable.
    """
    
    user_prompt = f"""
    Template sections: {template_text}
    
    Uploaded files summary:
    {files_text}
    
    What clarifying questions would help create a better document?
    """
    
    try:
        response = client.chat.completions.create(
            model="gpt-4o-mini",  # Use cheaper model for questions
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt},
            ]
        )
        return response.choices[0].message.content
    except Exception as e:
        print(f"Error calling OpenAI API for questions: {e}")
        return "What is the target audience for this document?\nWhat specific timeframe should the analysis cover?\nAre there any particular metrics or KPIs you'd like to focus on?"